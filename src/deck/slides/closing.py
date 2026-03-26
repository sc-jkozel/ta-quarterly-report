"""Slide 7: Closing statement."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from src.config import ReportConfig
from src.data.models import Insights
from src.deck.styles import (
    DARK_CHARCOAL, WHITE_RABBIT, PURPLE_RAIN, MID_GRAY,
    FONT_HEADING, FONT_BODY, SLIDE_WIDTH,
)
from src.deck.slides.helpers import add_blank_slide, set_slide_bg, add_textbox


def build(prs: Presentation, config: ReportConfig, insights: Insights, **kwargs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, DARK_CHARCOAL)

    # Accent bar
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE_RAIN
    bar.line.fill.background()

    # Headline
    add_textbox(
        slide, "Technical Architects make",
        Inches(0.75), Inches(1.2), Inches(11), Inches(0.5),
        font_name=FONT_HEADING, font_size=Pt(28),
        font_color=WHITE_RABBIT, bold=True,
    )
    add_textbox(
        slide, "complex deals closeable.",
        Inches(0.75), Inches(1.7), Inches(11), Inches(0.5),
        font_name=FONT_HEADING, font_size=Pt(28),
        font_color=PURPLE_RAIN, bold=True,
    )

    # Closing statement
    add_textbox(
        slide, insights.closing_statement,
        Inches(0.75), Inches(3.0), Inches(10), Inches(1.5),
        font_name=FONT_BODY, font_size=Pt(13),
        font_color=MID_GRAY,
    )

    # Attribution
    add_textbox(
        slide,
        f"{config.author_name} | {config.author_title} | {config.author_email}",
        Inches(0.75), Inches(6.5), Inches(10), Inches(0.3),
        font_name=FONT_BODY, font_size=Pt(10),
        font_color=MID_GRAY,
    )
