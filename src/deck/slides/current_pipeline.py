"""Slide 3: Current TA Pipeline - summary cards + stage breakdown chart."""

from collections import Counter

from pptx.presentation import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from src.deck.styles import (
    DARK_CHARCOAL, WHITE_RABBIT, PURPLE_RAIN, BLUE_MOON,
    YELLOW_SUBMARINE, MID_GRAY, CARD_BG_DARK,
    FONT_HEADING, FONT_BODY,
)
from src.deck.slides.helpers import (
    add_blank_slide, set_slide_bg, add_textbox, add_rounded_rect,
)
from src.data.models import ReportData
from src.insights.metrics import format_currency

# One accent color per stage — visually progresses from cool to warm.
STAGE_COLORS = [
    RGBColor(0x7C, 0xB3, 0xFF),        # Sky Blue   (Qualify)
    BLUE_MOON,                          # #00D1FF    (Discovery)
    PURPLE_RAIN,                        # #6559FF    (Pre-Sales)
    YELLOW_SUBMARINE,                   # #FFD700    (Proposal)
    RGBColor(0x00, 0xC8, 0x53),         # Green      (Negotiation)
    RGBColor(0xFF, 0xA0, 0x00),         # Amber      (overflow)
    RGBColor(0xE0, 0x6C, 0xFF),         # Violet     (overflow)
    RGBColor(0x00, 0xE6, 0xA0),         # Mint       (overflow)
]

# Pipeline stage order (early → late)
STAGE_ORDER = ["Qualify", "Discovery", "Pre-Sales", "Proposal", "Negotiation"]


def build(prs: Presentation, data: ReportData, metrics: dict, **kwargs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, DARK_CHARCOAL)

    # ── Title ──
    add_textbox(
        slide, "Current TA Pipeline",
        Inches(0.75), Inches(0.4), Inches(6), Inches(0.5),
        font_name=FONT_HEADING, font_size=Pt(24),
        font_color=WHITE_RABBIT, bold=True,
    )

    # ── Summary cards ──
    open_opps = [
        o for o in data.raw_opportunities
        if o.stage not in ("Closed Won", "Closed Lost")
    ]
    open_aud = sum(o.rollup_amount_aud for o in open_opps)
    open_count = len(open_opps)
    avg_pipeline_deal = open_aud / open_count if open_count > 0 else 0

    cards = [
        (format_currency(metrics["open_pipeline_aud"]), "Open Pipeline (AUD)"),
        (format_currency(avg_pipeline_deal), "Avg Deal Size in Pipeline"),
        (str(open_count), "Opportunities in Pipeline"),
    ]

    card_w = Inches(3.75)
    card_h = Inches(1.3)
    start_x = Inches(0.75)
    gap = Inches(0.29)
    cards_y = Inches(1.15)

    for i, (val, label) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rounded_rect(slide, x, cards_y, card_w, card_h, fill_color=CARD_BG_DARK)
        add_textbox(
            slide, val,
            x + Inches(0.3), cards_y + Inches(0.2),
            card_w - Inches(0.6), Inches(0.5),
            font_name=FONT_HEADING, font_size=Pt(28),
            font_color=BLUE_MOON, bold=True,
        )
        add_textbox(
            slide, label,
            x + Inches(0.3), cards_y + Inches(0.8),
            card_w - Inches(0.6), Inches(0.3),
            font_name=FONT_BODY, font_size=Pt(10),
            font_color=MID_GRAY,
        )

    # ── Stage breakdown chart ──
    chart_title_y = Inches(2.75)
    add_textbox(
        slide, "Opportunities by Stage",
        Inches(0.75), chart_title_y, Inches(6), Inches(0.4),
        font_name=FONT_BODY, font_size=Pt(14),
        font_color=WHITE_RABBIT, bold=True,
    )

    # Count opportunities per stage (exclude closed)
    stage_counts = Counter()
    for opp in data.raw_opportunities:
        if opp.stage not in ("Closed Won", "Closed Lost"):
            stage_counts[opp.stage] += 1

    if not stage_counts:
        add_textbox(
            slide, "No open pipeline data available",
            Inches(0.75), Inches(3.5), Inches(6), Inches(0.5),
            font_name=FONT_BODY, font_size=Pt(12),
            font_color=MID_GRAY,
        )
        return

    sorted_stages = [
        (stage, stage_counts[stage])
        for stage in STAGE_ORDER
        if stage in stage_counts
    ]
    # Append any unexpected stages at the end
    for stage, count in stage_counts.most_common():
        if stage not in STAGE_ORDER:
            sorted_stages.append((stage, count))

    max_count = max(count for _, count in sorted_stages)

    # Chart geometry — left-aligned with cards
    label_left = Inches(0.75)
    label_width = Inches(1.8)
    chart_left = Inches(2.7)
    max_bar_right = Inches(10.5)
    max_bar_width = max_bar_right - chart_left
    chart_start_y = Inches(3.35)
    chart_bottom = Inches(7.0)

    num_rows = len(sorted_stages)
    available = chart_bottom - chart_start_y
    row_height = int(available / num_rows)
    row_height = min(row_height, Inches(0.7))
    bar_height = int(row_height * 0.55)
    bar_pad = int((row_height - bar_height) / 2)

    for i, (stage, count) in enumerate(sorted_stages):
        y = chart_start_y + i * row_height
        color = STAGE_COLORS[i % len(STAGE_COLORS)]

        # Stage label (right-aligned, vertically centered to bar)
        tb = add_textbox(
            slide, stage,
            label_left, y, label_width, row_height,
            font_name=FONT_BODY, font_size=Pt(11),
            font_color=WHITE_RABBIT,
            alignment=PP_ALIGN.RIGHT,
        )
        tb.text_frame.paragraphs[0].space_before = Pt(0)
        tb.text_frame.paragraphs[0].space_after = Pt(0)
        tb.text_frame.auto_size = None
        tb.text_frame.word_wrap = True
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
        # Vertically center text in textbox
        txBody = tb._element.txBody
        bodyPr = txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        bodyPr.set('anchor', 'ctr')

        # Bar
        bar_ratio = count / max_count if max_count > 0 else 0
        bar_width = int(max_bar_width * bar_ratio)
        bar_width = max(bar_width, Inches(0.4))

        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            chart_left, y + bar_pad,
            bar_width, bar_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.adjustments[0] = 0.15

        # Count label (vertically centered to bar)
        count_tb = add_textbox(
            slide, str(count),
            chart_left + bar_width + Inches(0.15), y,
            Inches(1.5), row_height,
            font_name=FONT_BODY, font_size=Pt(11),
            font_color=MID_GRAY, bold=True,
        )
        count_body = count_tb._element.txBody
        count_bodyPr = count_body.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        count_bodyPr.set('anchor', 'ctr')
