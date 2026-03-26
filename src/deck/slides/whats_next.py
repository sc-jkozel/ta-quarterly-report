"""Slide 6: What's Next - numbered action items."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from src.data.models import Insights
from src.deck.styles import (
    WHITE_RABBIT, DARK_CHARCOAL, PURPLE_RAIN, MID_GRAY,
    FONT_HEADING, FONT_BODY,
)
from src.deck.slides.helpers import (
    add_blank_slide, set_slide_bg, add_textbox, add_rounded_rect,
)


def build(prs: Presentation, insights: Insights, **kwargs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, DARK_CHARCOAL)

    # Title
    add_textbox(
        slide, "What\u2019s Next",
        Inches(0.75), Inches(0.4), Inches(6), Inches(0.5),
        font_name=FONT_HEADING, font_size=Pt(24),
        font_color=WHITE_RABBIT, bold=True,
    )

    # Action items
    item_height = Inches(1.2)
    start_y = Inches(1.4)
    gap = Inches(0.25)

    for i, item in enumerate(insights.whats_next):
        y = start_y + i * (item_height + gap)

        # Number circle
        from pptx.enum.shapes import MSO_SHAPE
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.75), y + Inches(0.1),
            Inches(0.55), Inches(0.55),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = PURPLE_RAIN
        circle.line.fill.background()

        # Number text
        tf = circle.text_frame
        tf.word_wrap = False
        from pptx.enum.text import MSO_ANCHOR
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = f"{i + 1:02d}"
        run.font.name = FONT_HEADING
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE_RABBIT
        run.font.bold = True

        # Heading
        add_textbox(
            slide, item.heading,
            Inches(1.6), y + Inches(0.05),
            Inches(10), Inches(0.35),
            font_name=FONT_BODY, font_size=Pt(14),
            font_color=WHITE_RABBIT, bold=True,
        )

        # Description
        add_textbox(
            slide, item.description,
            Inches(1.6), y + Inches(0.45),
            Inches(10), Inches(0.7),
            font_name=FONT_BODY, font_size=Pt(11),
            font_color=MID_GRAY,
        )
