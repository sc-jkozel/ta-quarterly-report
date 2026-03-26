"""Slide 1: Title slide."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from src.config import ReportConfig
from src.deck.styles import (
    DARK_CHARCOAL, WHITE_RABBIT, PURPLE_RAIN, MID_GRAY,
    FONT_HEADING, FONT_BODY, SLIDE_WIDTH,
)
from src.deck.slides.helpers import add_blank_slide, set_slide_bg, add_textbox


def build(prs: Presentation, config: ReportConfig, **kwargs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, DARK_CHARCOAL)

    # Accent bar
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE_RAIN
    bar.line.fill.background()

    # Title
    add_textbox(
        slide, "Technical Architect",
        Inches(0.75), Inches(1.5), Inches(11), Inches(0.6),
        font_name=FONT_HEADING, font_size=Pt(40),
        font_color=WHITE_RABBIT, bold=True,
        alignment=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide, "Impact Report",
        Inches(0.75), Inches(2.1), Inches(11), Inches(0.6),
        font_name=FONT_HEADING, font_size=Pt(40),
        font_color=PURPLE_RAIN, bold=True,
        alignment=PP_ALIGN.LEFT,
    )

    # Subtitle
    add_textbox(
        slide, f"{config.fiscal_year} Executive Readout",
        Inches(0.75), Inches(3.2), Inches(6), Inches(0.4),
        font_name=FONT_BODY, font_size=Pt(16),
        font_color=WHITE_RABBIT, bold=False,
    )
    add_textbox(
        slide, f"Prepared for {config.prepared_for}",
        Inches(0.75), Inches(3.7), Inches(6), Inches(0.35),
        font_name=FONT_BODY, font_size=Pt(12),
        font_color=MID_GRAY,
    )

    # Footer
    import datetime
    month_year = datetime.date.today().strftime("%B %Y")
    add_textbox(
        slide,
        f"{config.org_name} | {config.team_name} | {month_year}",
        Inches(0.75), Inches(6.5), Inches(10), Inches(0.3),
        font_name=FONT_BODY, font_size=Pt(10),
        font_color=MID_GRAY,
    )
