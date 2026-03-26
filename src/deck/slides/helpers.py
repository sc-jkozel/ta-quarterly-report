"""Shared slide building utilities."""

from __future__ import annotations

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from src.deck.styles import (
    DARK_CHARCOAL, WHITE_RABBIT, PURPLE_RAIN, BLUE_MOON, MID_GRAY,
    CARD_BG_DARK,
    FONT_HEADING, FONT_BODY, SIZE_HEADING, SIZE_SUBHEADING, SIZE_BODY,
)


def add_blank_slide(prs: Presentation) -> Slide:
    """Add a blank slide to the presentation."""
    layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(layout)


def set_slide_bg(slide: Slide, color: RGBColor):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide: Slide,
    text: str,
    left: Inches | Emu,
    top: Inches | Emu,
    width: Inches | Emu,
    height: Inches | Emu,
    font_name: str = FONT_BODY,
    font_size: Pt = SIZE_BODY,
    font_color: RGBColor = WHITE_RABBIT,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
):
    """Add a styled textbox with a single run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox


def add_rounded_rect(
    slide: Slide,
    left, top, width, height,
    fill_color: RGBColor,
    border_color: RGBColor | None = None,
    corner_radius: int = 80000,
):
    """Add a rounded rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Set corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_metric_card(
    slide: Slide,
    value: str,
    label: str,
    left, top, width, height,
    value_color: RGBColor = BLUE_MOON,
    bg_color: RGBColor = CARD_BG_DARK,
):
    """Add a metric display card with large value and small label."""
    rect = add_rounded_rect(slide, left, top, width, height, bg_color)

    # Value
    add_textbox(
        slide, value,
        left + Inches(0.2), top + Inches(0.2),
        width - Inches(0.4), Inches(0.6),
        font_name=FONT_HEADING, font_size=Pt(28),
        font_color=value_color, bold=True,
        alignment=PP_ALIGN.LEFT,
    )
    # Label
    add_textbox(
        slide, label,
        left + Inches(0.2), top + height - Inches(0.55),
        width - Inches(0.4), Inches(0.35),
        font_name=FONT_BODY, font_size=Pt(10),
        font_color=MID_GRAY,
        alignment=PP_ALIGN.LEFT,
    )
    return rect


def add_bullet_list(
    slide: Slide,
    items: list[str],
    left, top, width, height,
    font_color: RGBColor = WHITE_RABBIT,
    font_size: Pt = SIZE_BODY,
    bullet_char: str = "\u2022",
):
    """Add a bullet list as a textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.name = FONT_BODY
        run.font.size = font_size
        run.font.color.rgb = font_color
    return txBox
