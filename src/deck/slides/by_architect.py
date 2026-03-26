"""Slide 4: By Technical Architect - per-TA breakdown cards."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from src.config import TECHNICAL_ARCHITECTS
from src.deck.styles import (
    DARK_CHARCOAL, WHITE_RABBIT, PURPLE_RAIN, BLUE_MOON, MID_GRAY,
    CARD_BG_DARK, FONT_HEADING, FONT_BODY,
)
from src.deck.slides.helpers import (
    add_blank_slide, set_slide_bg, add_textbox, add_rounded_rect,
)


def build(prs: Presentation, ta_cards: list[dict], config=None, **kwargs):
    from src.config import ReportConfig
    if config is None:
        config = ReportConfig()

    slide = add_blank_slide(prs)
    set_slide_bg(slide, DARK_CHARCOAL)

    # Title
    add_textbox(
        slide, f"{config.fiscal_year.upper()} Performance by Architect",
        Inches(0.75), Inches(0.4), Inches(6), Inches(0.5),
        font_name=FONT_HEADING, font_size=Pt(24),
        font_color=WHITE_RABBIT, bold=True,
    )

    # Build TA info lookup
    ta_info = {ta["name"]: ta for ta in TECHNICAL_ARCHITECTS}

    # Layout: 2x2 grid for 4 TAs
    card_w = Inches(5.7)
    card_h = Inches(2.5)
    x_positions = [Inches(0.75), Inches(6.85)]
    y_positions = [Inches(1.3), Inches(4.1)]

    for i, card in enumerate(ta_cards[:4]):
        col = i % 2
        row = i // 2
        x = x_positions[col]
        y = y_positions[row]

        # Card background
        add_rounded_rect(slide, x, y, card_w, card_h, fill_color=CARD_BG_DARK)

        # Accent bar on left
        from pptx.enum.shapes import MSO_SHAPE
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y, Inches(0.06), card_h,
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = PURPLE_RAIN
        accent.line.fill.background()

        # Name + region
        info = ta_info.get(card["name"], {})
        region = info.get("region", "")
        location = info.get("location", "")
        region_text = f"{region} ({location})" if location else region

        add_textbox(
            slide, card["name"],
            x + Inches(0.3), y + Inches(0.15),
            Inches(4), Inches(0.35),
            font_name=FONT_HEADING, font_size=Pt(16),
            font_color=WHITE_RABBIT, bold=True,
        )
        add_textbox(
            slide, region_text,
            x + Inches(0.3), y + Inches(0.5),
            Inches(4), Inches(0.25),
            font_name=FONT_BODY, font_size=Pt(10),
            font_color=MID_GRAY,
        )

        # Metrics grid - 2 columns
        metrics_left = [
            ("Opportunities", str(card["opportunities"])),
            ("Open Pipeline (AUD)", card["open_pipeline_fmt"]),
        ]
        metrics_right = [
            ("Won / Lost", card["record"]),
            ("Booked (AUD)", card["booked_fmt"]),
        ]

        for j, (label, val) in enumerate(metrics_left):
            my = y + Inches(0.95) + j * Inches(0.55)
            add_textbox(
                slide, label,
                x + Inches(0.3), my,
                Inches(2.2), Inches(0.22),
                font_name=FONT_BODY, font_size=Pt(9),
                font_color=MID_GRAY,
            )
            add_textbox(
                slide, val,
                x + Inches(0.3), my + Inches(0.2),
                Inches(2.2), Inches(0.25),
                font_name=FONT_BODY, font_size=Pt(13),
                font_color=BLUE_MOON, bold=True,
            )

        for j, (label, val) in enumerate(metrics_right):
            my = y + Inches(0.95) + j * Inches(0.55)
            add_textbox(
                slide, label,
                x + Inches(2.8), my,
                Inches(2.2), Inches(0.22),
                font_name=FONT_BODY, font_size=Pt(9),
                font_color=MID_GRAY,
            )
            add_textbox(
                slide, val,
                x + Inches(2.8), my + Inches(0.2),
                Inches(2.2), Inches(0.25),
                font_name=FONT_BODY, font_size=Pt(13),
                font_color=BLUE_MOON, bold=True,
            )

    # Note for any newly-joined TAs not yet reflected in data
    new_tas = [
        ta for ta in TECHNICAL_ARCHITECTS if ta.get("new")
    ]
    if new_tas:
        names = ", ".join(
            f"{ta['name']} ({ta['region']})" for ta in new_tas
        )
        note = f"{names} recently joined the team — not yet reflected in opportunity data."
        add_textbox(
            slide, note,
            Inches(0.75), Inches(6.8), Inches(10), Inches(0.3),
            font_name=FONT_BODY, font_size=Pt(9),
            font_color=MID_GRAY, alignment=PP_ALIGN.LEFT,
        )
