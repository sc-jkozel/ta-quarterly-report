"""Slide 5: Reading the Data - Strong Signals / Watch & Evolve."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from src.data.models import Insights
from src.deck.styles import (
    WHITE_RABBIT, DARK_CHARCOAL, PURPLE_RAIN, BLUE_MOON,
    YELLOW_SUBMARINE, MID_GRAY, CARD_BG_DARK,
    FONT_HEADING, FONT_BODY,
)
from src.deck.slides.helpers import (
    add_blank_slide, set_slide_bg, add_textbox, add_rounded_rect,
    add_bullet_list,
)


def build(prs: Presentation, insights: Insights, **kwargs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, DARK_CHARCOAL)

    # Title
    add_textbox(
        slide, "Reading the Data",
        Inches(0.75), Inches(0.4), Inches(6), Inches(0.5),
        font_name=FONT_HEADING, font_size=Pt(24),
        font_color=WHITE_RABBIT, bold=True,
    )
    add_textbox(
        slide, "What the numbers tell us \u2014 and where to be careful",
        Inches(0.75), Inches(0.95), Inches(8), Inches(0.3),
        font_name=FONT_BODY, font_size=Pt(12),
        font_color=MID_GRAY,
    )

    col_w = Inches(5.7)
    col_h = Inches(5.0)
    y = Inches(1.6)

    # Left column - Strong Signals
    left_x = Inches(0.75)
    add_rounded_rect(
        slide, left_x, y, col_w, col_h,
        fill_color=CARD_BG_DARK,
    )

    # Accent bar
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left_x, y, col_w, Inches(0.05),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_MOON
    bar.line.fill.background()

    add_textbox(
        slide, "Strong Signals",
        left_x + Inches(0.3), y + Inches(0.2),
        Inches(4), Inches(0.35),
        font_name=FONT_BODY, font_size=Pt(14),
        font_color=BLUE_MOON, bold=True,
    )

    add_bullet_list(
        slide, insights.strong_signals,
        left_x + Inches(0.3), y + Inches(0.7),
        col_w - Inches(0.6), col_h - Inches(1.0),
        font_color=WHITE_RABBIT, font_size=Pt(11),
    )

    # Right column - Watch & Evolve
    right_x = Inches(6.85)
    add_rounded_rect(
        slide, right_x, y, col_w, col_h,
        fill_color=CARD_BG_DARK,
    )

    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        right_x, y, col_w, Inches(0.05),
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = YELLOW_SUBMARINE
    bar2.line.fill.background()

    add_textbox(
        slide, "Watch & Evolve",
        right_x + Inches(0.3), y + Inches(0.2),
        Inches(4), Inches(0.35),
        font_name=FONT_BODY, font_size=Pt(14),
        font_color=YELLOW_SUBMARINE, bold=True,
    )

    add_bullet_list(
        slide, insights.watch_evolve,
        right_x + Inches(0.3), y + Inches(0.7),
        col_w - Inches(0.6), col_h - Inches(1.0),
        font_color=WHITE_RABBIT, font_size=Pt(11),
    )
